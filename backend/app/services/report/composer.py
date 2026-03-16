"""
보고서 생성기 — Jinja2 → HTML → PDF + python-pptx
"""

import asyncio
import os
from typing import Optional
from jinja2 import Environment, FileSystemLoader
import structlog

from app.services.snowflake.client import get_snowflake_client

log = structlog.get_logger(__name__)

TEMPLATE_DIR = os.path.join(os.path.dirname(__file__), "templates")
jinja_env = Environment(loader=FileSystemLoader(TEMPLATE_DIR), autoescape=True)


async def compose_report(req: dict, job_id: str) -> str:
    """
    보고서 유형에 따라 데이터 수집 → 템플릿 렌더 → PDF/PPT 생성
    Returns: 생성된 파일 URL
    """
    report_type = req["report_type"]
    dept        = req["department_code"]

    log.info("Report composition started", job_id=job_id, type=report_type, dept=dept)

    # 1) 데이터 수집
    data = await _collect_report_data(report_type, dept, req["period_start"], req["period_end"])

    # 2) 포맷별 생성
    fmt = req.get("format", "pdf")
    if fmt == "pdf":
        url = await _render_pdf(report_type, dept, data, job_id)
    elif fmt == "pptx":
        url = await _render_pptx(report_type, dept, data, job_id)
    else:
        url = await _render_pdf(report_type, dept, data, job_id)

    log.info("Report composition done", job_id=job_id, url=url)
    return url


async def _collect_report_data(report_type: str, dept: str, period_start, period_end) -> dict:
    client = get_snowflake_client()
    data: dict = {"dept": dept, "period_start": str(period_start), "period_end": str(period_end)}

    try:
        if report_type in ("weekly_ops", "monthly_mgmt"):
            data["kpis"] = await client.execute("""
                SELECT kpi_name, kpi_value, kpi_unit, trend_pct, status
                FROM FACT_DB.MART.mart_exec_kpi_weekly
                WHERE week_label = TO_CHAR(%(start)s::DATE, 'IYYY-IW')
            """, {"start": str(period_start)})

        if dept in ("production", "all"):
            data["oee"] = await client.execute("""
                SELECT production_date, line_code, oee, total_downtime_h
                FROM FACT_DB.MART.mart_production_oee_daily
                WHERE production_date BETWEEN %(start)s AND %(end)s
                ORDER BY production_date
            """, {"start": str(period_start), "end": str(period_end)})
    except Exception as e:
        log.warning("Data collection partial failure", error=str(e))

    return data


async def _render_pdf(report_type: str, dept: str, data: dict, job_id: str) -> str:
    """Jinja2 → HTML → WeasyPrint PDF"""
    try:
        template = jinja_env.get_template(f"{report_type}.html")
        html = template.render(**data)
    except Exception:
        html = f"<h1>F.A.C.T 보고서</h1><pre>{data}</pre>"

    output_path = f"/tmp/report_{job_id}.pdf"

    def _convert():
        try:
            from weasyprint import HTML
            HTML(string=html).write_pdf(output_path)
        except ImportError:
            with open(output_path, "wb") as f:
                f.write(b"%PDF-1.4 placeholder")

    await asyncio.to_thread(_convert)
    # TODO: upload to S3/Azure Blob and return URL
    return f"s3://fact-reports/{job_id}.pdf"


async def _render_pptx(report_type: str, dept: str, data: dict, job_id: str) -> str:
    """python-pptx 기반 PPT 생성"""
    output_path = f"/tmp/report_{job_id}.pptx"

    def _build():
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = f"F.A.C.T 보고서 — {dept} ({data['period_start']}~{data['period_end']})"

        prs.save(output_path)

    await asyncio.to_thread(_build)
    return f"s3://fact-reports/{job_id}.pptx"
